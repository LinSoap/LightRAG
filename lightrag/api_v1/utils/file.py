from pathlib import Path
from pdb import pm
import traceback

import aiofiles
from fastapi import logger

from lightrag.lightrag import LightRAG

# Temporary file prefix
temp_prefix = "__tmp__"


def sanitize_filename(filename: str, input_dir: Path) -> str:
    """
    Sanitize uploaded filename to prevent Path Traversal attacks.

    Args:
        filename: The original filename from the upload
        input_dir: The target input directory

    Returns:
        str: Sanitized filename that is safe to use

    Raises:
        ValueError: If the filename is unsafe or invalid
    """
    # Basic validation
    if not filename or not filename.strip():
        raise ValueError("Filename cannot be empty")

    # Remove path separators and traversal sequences
    clean_name = filename.replace("/", "").replace("\\", "")
    clean_name = clean_name.replace("..", "")

    # Remove control characters and null bytes
    clean_name = "".join(c for c in clean_name if ord(c) >= 32 and c != "\x7f")

    # Remove leading/trailing whitespace and dots
    clean_name = clean_name.strip().strip(".")

    # Check if anything is left after sanitization
    if not clean_name:
        raise ValueError("Invalid filename after sanitization")

    # Verify the final path stays within the input directory
    try:
        final_path = (input_dir / clean_name).resolve()
        if not final_path.is_relative_to(input_dir.resolve()):
            raise ValueError("Unsafe filename detected")
    except (OSError, ValueError):
        raise ValueError("Invalid filename")

    return clean_name


async def pipeline_enqueue_file(
    rag: LightRAG, file_path: Path, track_id: str = None
) -> tuple[bool, str]:
    """Add a file to the queue for processing

    Args:
        rag: LightRAG instance
        file_path: Path to the saved file
        track_id: Optional tracking ID, if not provided will be generated
    Returns:
        tuple: (success: bool, track_id: str)
    """

    # Generate track_id if not provided
    if track_id is None:
        track_id = generate_track_id("unknown")

    try:
        content = ""
        ext = file_path.suffix.lower()
        file_size = 0

        # Get file size for error reporting
        try:
            file_size = file_path.stat().st_size
        except Exception:
            file_size = 0

        file = None
        try:
            async with aiofiles.open(file_path, "rb") as f:
                file = await f.read()
        except PermissionError as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]Permission denied - cannot read file",
                    "original_error": str(e),
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(
                f"[File Extraction]Permission denied reading file: {file_path.name}"
            )
            return False, track_id
        except FileNotFoundError as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]File not found",
                    "original_error": str(e),
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(f"[File Extraction]File not found: {file_path.name}")
            return False, track_id
        except Exception as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]File reading error",
                    "original_error": str(e),
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(
                f"[File Extraction]Error reading file {file_path.name}: {str(e)}"
            )
            return False, track_id

        # Process based on file type
        try:
            match ext:
                case (
                    ".txt"
                    | ".md"
                    | ".html"
                    | ".htm"
                    | ".tex"
                    | ".json"
                    | ".xml"
                    | ".yaml"
                    | ".yml"
                    | ".rtf"
                    | ".odt"
                    | ".epub"
                    | ".csv"
                    | ".log"
                    | ".conf"
                    | ".ini"
                    | ".properties"
                    | ".sql"
                    | ".bat"
                    | ".sh"
                    | ".c"
                    | ".cpp"
                    | ".py"
                    | ".java"
                    | ".js"
                    | ".ts"
                    | ".swift"
                    | ".go"
                    | ".rb"
                    | ".php"
                    | ".css"
                    | ".scss"
                    | ".less"
                ):
                    try:
                        # Try to decode as UTF-8
                        content = file.decode("utf-8")

                        # Validate content
                        if not content or len(content.strip()) == 0:
                            error_files = [
                                {
                                    "file_path": str(file_path.name),
                                    "error_description": "[File Extraction]Empty file content",
                                    "original_error": "File contains no content or only whitespace",
                                    "file_size": file_size,
                                }
                            ]
                            await rag.apipeline_enqueue_error_documents(
                                error_files, track_id
                            )
                            logger.error(
                                f"[File Extraction]Empty content in file: {file_path.name}"
                            )
                            return False, track_id

                        # Check if content looks like binary data string representation
                        if content.startswith("b'") or content.startswith('b"'):
                            error_files = [
                                {
                                    "file_path": str(file_path.name),
                                    "error_description": "[File Extraction]Binary data in text file",
                                    "original_error": "File appears to contain binary data representation instead of text",
                                    "file_size": file_size,
                                }
                            ]
                            await rag.apipeline_enqueue_error_documents(
                                error_files, track_id
                            )
                            logger.error(
                                f"[File Extraction]File {file_path.name} appears to contain binary data representation instead of text"
                            )
                            return False, track_id

                    except UnicodeDecodeError as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]UTF-8 encoding error, please convert it to UTF-8 before processing",
                                "original_error": f"File is not valid UTF-8 encoded text: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]File {file_path.name} is not valid UTF-8 encoded text. Please convert it to UTF-8 before processing."
                        )
                        return False, track_id

                case ".pdf":
                    try:
                        if not pm.is_installed("pypdf2"):  # type: ignore
                            pm.install("pypdf2")
                        from PyPDF2 import PdfReader  # type: ignore
                        from io import BytesIO

                        pdf_file = BytesIO(file)
                        reader = PdfReader(pdf_file)
                        for page in reader.pages:
                            content += page.extract_text() + "\n"
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]PDF processing error",
                                "original_error": f"Failed to extract text from PDF: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing PDF {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case ".docx":
                    try:
                        if not pm.is_installed("python-docx"):  # type: ignore
                            try:
                                pm.install("python-docx")
                            except Exception:
                                pm.install("docx")
                        from docx import Document  # type: ignore
                        from io import BytesIO

                        docx_file = BytesIO(file)
                        doc = Document(docx_file)
                        content = "\n".join(
                            [paragraph.text for paragraph in doc.paragraphs]
                        )
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]DOCX processing error",
                                "original_error": f"Failed to extract text from DOCX: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing DOCX {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case ".pptx":
                    try:
                        if not pm.is_installed("python-pptx"):  # type: ignore
                            pm.install("pptx")
                        from pptx import Presentation  # type: ignore
                        from io import BytesIO

                        pptx_file = BytesIO(file)
                        prs = Presentation(pptx_file)
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    content += shape.text + "\n"
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]PPTX processing error",
                                "original_error": f"Failed to extract text from PPTX: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing PPTX {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case ".xlsx":
                    try:
                        if not pm.is_installed("openpyxl"):  # type: ignore
                            pm.install("openpyxl")
                        from openpyxl import load_workbook  # type: ignore
                        from io import BytesIO

                        xlsx_file = BytesIO(file)
                        wb = load_workbook(xlsx_file)
                        for sheet in wb:
                            content += f"Sheet: {sheet.title}\n"
                            for row in sheet.iter_rows(values_only=True):
                                content += (
                                    "\t".join(
                                        str(cell) if cell is not None else ""
                                        for cell in row
                                    )
                                    + "\n"
                                )
                            content += "\n"
                    except Exception as e:
                        error_files = [
                            {
                                "file_path": str(file_path.name),
                                "error_description": "[File Extraction]XLSX processing error",
                                "original_error": f"Failed to extract text from XLSX: {str(e)}",
                                "file_size": file_size,
                            }
                        ]
                        await rag.apipeline_enqueue_error_documents(
                            error_files, track_id
                        )
                        logger.error(
                            f"[File Extraction]Error processing XLSX {file_path.name}: {str(e)}"
                        )
                        return False, track_id

                case _:
                    error_files = [
                        {
                            "file_path": str(file_path.name),
                            "error_description": f"[File Extraction]Unsupported file type: {ext}",
                            "original_error": f"File extension {ext} is not supported",
                            "file_size": file_size,
                        }
                    ]
                    await rag.apipeline_enqueue_error_documents(error_files, track_id)
                    logger.error(
                        f"[File Extraction]Unsupported file type: {file_path.name} (extension {ext})"
                    )
                    return False, track_id

        except Exception as e:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "[File Extraction]File format processing error",
                    "original_error": f"Unexpected error during file extracting: {str(e)}",
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(
                f"[File Extraction]Unexpected error during {file_path.name} extracting: {str(e)}"
            )
            return False, track_id

        # Insert into the RAG queue
        if content:
            # Check if content contains only whitespace characters
            if not content.strip():
                error_files = [
                    {
                        "file_path": str(file_path.name),
                        "error_description": "[File Extraction]File contains only whitespace",
                        "original_error": "File content contains only whitespace characters",
                        "file_size": file_size,
                    }
                ]
                await rag.apipeline_enqueue_error_documents(error_files, track_id)
                logger.warning(
                    f"[File Extraction]File contains only whitespace characters: {file_path.name}"
                )
                return False, track_id

            try:
                await rag.apipeline_enqueue_documents(
                    content, file_paths=file_path.name, track_id=track_id
                )

                logger.info(
                    f"Successfully extracted and enqueued file: {file_path.name}"
                )

                # Move file to __enqueued__ directory after enqueuing
                try:
                    enqueued_dir = file_path.parent / "__enqueued__"
                    enqueued_dir.mkdir(exist_ok=True)

                    # Generate unique filename to avoid conflicts
                    unique_filename = get_unique_filename_in_enqueued(
                        enqueued_dir, file_path.name
                    )
                    target_path = enqueued_dir / unique_filename

                    # Move the file
                    file_path.rename(target_path)
                    logger.debug(
                        f"Moved file to enqueued directory: {file_path.name} -> {unique_filename}"
                    )

                except Exception as move_error:
                    pass
                    logger.error(
                        f"Failed to move file {file_path.name} to __enqueued__ directory: {move_error}"
                    )
                    # Don't affect the main function's success status

                return True, track_id

            except Exception as e:
                error_files = [
                    {
                        "file_path": str(file_path.name),
                        "error_description": "Document enqueue error",
                        "original_error": f"Failed to enqueue document: {str(e)}",
                        "file_size": file_size,
                    }
                ]
                await rag.apipeline_enqueue_error_documents(error_files, track_id)
                logger.error(f"Error enqueueing document {file_path.name}: {str(e)}")
                return False, track_id
        else:
            error_files = [
                {
                    "file_path": str(file_path.name),
                    "error_description": "No content extracted",
                    "original_error": "No content could be extracted from file",
                    "file_size": file_size,
                }
            ]
            await rag.apipeline_enqueue_error_documents(error_files, track_id)
            logger.error(f"No content extracted from file: {file_path.name}")
            return False, track_id
    except Exception as e:
        # Catch-all for any unexpected errors
        try:
            file_size = file_path.stat().st_size if file_path.exists() else 0
        except Exception:
            file_size = 0

        error_files = [
            {
                "file_path": str(file_path.name),
                "error_description": "Unexpected processing error",
                "original_error": f"Unexpected error: {str(e)}",
                "file_size": file_size,
            }
        ]
        await rag.apipeline_enqueue_error_documents(error_files, track_id)
        logger.error(f"Enqueuing file {file_path.name} error: {str(e)}")
        logger.error(traceback.format_exc())
        return False, track_id
    finally:
        if file_path.name.startswith(temp_prefix):
            try:
                file_path.unlink()
            except Exception as e:
                pass
                logger.error(f"Error deleting file {file_path}: {str(e)}")


def get_unique_filename_in_enqueued(target_dir: Path, original_name: str) -> str:
    """Generate a unique filename in the target directory by adding numeric suffixes if needed

    Args:
        target_dir: Target directory path
        original_name: Original filename

    Returns:
        str: Unique filename (may have numeric suffix added)
    """
    from pathlib import Path
    import time

    original_path = Path(original_name)
    base_name = original_path.stem
    extension = original_path.suffix

    # Try original name first
    if not (target_dir / original_name).exists():
        return original_name

    # Try with numeric suffixes 001-999
    for i in range(1, 1000):
        suffix = f"{i:03d}"
        new_name = f"{base_name}_{suffix}{extension}"
        if not (target_dir / new_name).exists():
            return new_name

    # Fallback with timestamp if all 999 slots are taken
    timestamp = int(time.time())
    return f"{base_name}_{timestamp}{extension}"
